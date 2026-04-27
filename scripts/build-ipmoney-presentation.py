from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "architecture" / "ipmoney-client-briefing.pptx"


NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(30, 64, 175)
LIGHT = RGBColor(248, 250, 252)
ACCENT = RGBColor(14, 116, 144)
TEXT = RGBColor(31, 41, 55)
WHITE = RGBColor(255, 255, 255)


def add_bg(slide, prs, color=LIGHT):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)


def add_title(slide, prs, title, subtitle=None):
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(10.8), Inches(0.55))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.65), Inches(1.05), Inches(12), Inches(0.6))
        stf = sbox.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.name = "Microsoft YaHei"
        srun.font.size = Pt(16)
        srun.font.color.rgb = BLUE


def add_bullets(slide, left, top, width, height, bullets, size=22):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.text = b
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)


def add_two_col_cards(slide, items):
    start_x = Inches(0.8)
    start_y = Inches(1.8)
    col_w = Inches(6.15)
    card_h = Inches(1.8)
    gap_x = Inches(0.25)
    gap_y = Inches(0.25)
    for idx, (title, body) in enumerate(items):
        r = idx // 2
        c = idx % 2
        x = start_x + c * (col_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        card = slide.shapes.add_shape(1, x, y, col_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(203, 213, 225)

        t = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.18), col_w - Inches(0.5), Inches(0.45))
        tf = t.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = "Microsoft YaHei"
        run.font.bold = True
        run.font.size = Pt(18)
        run.font.color.rgb = BLUE

        b = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.62), col_w - Inches(0.5), card_h - Inches(0.75))
        btf = b.text_frame
        btf.word_wrap = True
        btf.clear()
        bp = btf.paragraphs[0]
        brun = bp.add_run()
        brun.text = body
        brun.font.name = "Microsoft YaHei"
        brun.font.size = Pt(14)
        brun.font.color.rgb = TEXT


def add_footer(slide, text="Ipmoney | 甲方汇报版"):
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(226, 232, 240)
    line.line.fill.background()

    f = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(8), Inches(0.28))
    tf = f.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(10)
    r.font.color.rgb = RGBColor(100, 116, 139)

    d = slide.shapes.add_textbox(Inches(10.5), Inches(7.08), Inches(2.5), Inches(0.28))
    dtf = d.text_frame
    dtf.clear()
    dp = dtf.paragraphs[0]
    dp.alignment = PP_ALIGN.RIGHT
    dr = dp.add_run()
    dr.text = str(date.today())
    dr.font.name = "Microsoft YaHei"
    dr.font.size = Pt(10)
    dr.font.color.rgb = RGBColor(100, 116, 139)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, RGBColor(239, 246, 255))
    add_title(slide, prs, "Ipmoney 平台汇报", "知识产权交易与运营一体化平台（甲方宣讲版）")
    add_bullets(
        slide,
        Inches(0.9),
        Inches(2.0),
        Inches(8.5),
        Inches(3.5),
        [
            "聚焦四大重点：平台功能、核心优势、总体架构、应用场景",
            "支撑目标：经营可复制、流程可管理、作业可审计",
            "汇报对象：甲方管理层、业务负责人、运营与交付团队",
        ],
        size=22,
    )
    hero = slide.shapes.add_shape(1, Inches(9.3), Inches(1.8), Inches(3.2), Inches(3.8))
    hero.fill.solid()
    hero.fill.fore_color.rgb = BLUE
    hero.line.fill.background()
    htext = slide.shapes.add_textbox(Inches(9.55), Inches(2.2), Inches(2.8), Inches(3.0))
    htf = htext.text_frame
    htf.clear()
    for i, t in enumerate(["交易", "履约", "结算", "治理"]):
        p = htf.paragraphs[0] if i == 0 else htf.add_paragraph()
        p.text = t
        p.font.name = "Microsoft YaHei"
        p.font.bold = True
        p.font.size = Pt(30)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    add_footer(slide)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, prs, "1. 项目定位与业务价值")
    add_two_col_cards(
        slide,
        [
            ("对甲方", "形成可复制、可管理、可审计的平台经营能力，支撑规模化扩张。"),
            ("对用户", "降低交易门槛，提升交易透明度与履约确定性，增强成交信心。"),
            ("对运营团队", "建立标准作业流程与跨角色协同机制，提升人效与质量稳定性。"),
            ("平台主张", "以“找标的-谈交易-履约交付-资金结算-运营治理”构建统一闭环。"),
        ],
    )
    add_footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, prs, "2. 平台功能全景")
    add_two_col_cards(
        slide,
        [
            ("供需撮合", "搜索检索、专题运营、咨询会话、发布上架，提升线索到商机效率。"),
            ("交易履约", "订单、订金/尾款、合同、发票、退款、结算，全链路可追踪。"),
            ("运营治理", "审核、工单、会话分配、批量运营、配置中心，支撑标准化作业。"),
            ("风控合规", "RBAC、审计日志、幂等保护、告警闭环，保障资金与流程安全。"),
        ],
    )
    add_footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, prs, "3. 用户端核心能力")
    add_bullets(
        slide,
        Inches(0.9),
        Inches(1.9),
        Inches(12),
        Inches(4.8),
        [
            "首页与发现：多维检索（挂牌/成果/机构/发明人/技术经理）与运营推荐联动",
            "咨询与会话：从详情页一键发起咨询，支持持续多轮沟通",
            "发布与管理：支持专利交易/成果发布、认领申请、进度查看",
            "交易与履约：订单、分阶段支付、合同上传、发票下载、退款申请",
            "个人中心：认证、资料、通知、地址、客服和年费托管统一入口",
        ],
        size=20,
    )
    add_footer(slide)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, prs, "4. 管理端核心能力")
    add_two_col_cards(
        slide,
        [
            ("运营中心", "认证审核、挂牌审核、专利导入、批量上架、首页配置。"),
            ("客服中心", "会话中心、工单处理、争议升级、SLA 跟踪。"),
            ("财务中心", "退款审核、结算放款、发票归档、财务报表导出。"),
            ("治理中心", "账号角色管理、权限分配、审计追溯、告警处置。"),
        ],
    )
    add_footer(slide)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, prs, "5. 平台核心优势")
    add_bullets(
        slide,
        Inches(0.9),
        Inches(1.9),
        Inches(12.1),
        Inches(4.9),
        [
            "闭环优势：从线索到结算一体化，避免系统断点与流程断层",
            "协同优势：运营、客服、财务、管理员跨角色协同，职责清晰",
            "治理优势：查询-处理-留痕-复核标准作业，结果可追踪可复盘",
            "技术优势：统一 API 治理口径，模块化架构便于演进与扩展",
            "合规优势：身份门禁 + 权限门禁 + 幂等机制 + 审计快照",
        ],
        size=20,
    )
    add_footer(slide)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, prs, "6. 总体架构设计")
    layers = [
        ("用户层", "微信小程序为主，H5 为补充入口"),
        ("管理层", "运营/客服/财务/管理员统一后台作业"),
        ("业务服务层", "NestJS 模块化 API，按业务域组织能力"),
        ("数据与基础设施层", "PostgreSQL + Redis + 对象存储 + 审计与告警"),
    ]
    y = Inches(1.8)
    for title, body in layers:
        rect = slide.shapes.add_shape(1, Inches(0.9), y, Inches(11.8), Inches(1.15))
        rect.fill.solid()
        rect.fill.fore_color.rgb = WHITE
        rect.line.color.rgb = RGBColor(203, 213, 225)

        left = slide.shapes.add_shape(1, Inches(0.9), y, Inches(2.1), Inches(1.15))
        left.fill.solid()
        left.fill.fore_color.rgb = ACCENT
        left.line.fill.background()

        t1 = slide.shapes.add_textbox(Inches(1.05), y + Inches(0.32), Inches(1.8), Inches(0.5))
        tf1 = t1.text_frame
        tf1.clear()
        p1 = tf1.paragraphs[0]
        p1.text = title
        p1.font.name = "Microsoft YaHei"
        p1.font.bold = True
        p1.font.size = Pt(16)
        p1.font.color.rgb = WHITE

        t2 = slide.shapes.add_textbox(Inches(3.2), y + Inches(0.3), Inches(9.1), Inches(0.55))
        tf2 = t2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.text = body
        p2.font.name = "Microsoft YaHei"
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT

        y += Inches(1.28)
    add_footer(slide)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, prs, "7. 典型应用场景")
    add_two_col_cards(
        slide,
        [
            ("场景一：专利交易全流程", "检索咨询 -> 下单支付 -> 合同履约 -> 结算开票，支撑主交易场景。"),
            ("场景二：成果发布与认领", "发布审核 -> 上线展示 -> 认领审批 -> 归属回写，支撑供给侧运营。"),
            ("场景三：年费托管服务", "计划创建 -> 任务执行 -> 回执对账 -> 完成归档，形成增值服务闭环。"),
            ("场景四：争议与退款处理", "会话识别 -> 工单升级 -> 证据留存 -> 状态闭环，降低风险与投诉成本。"),
        ],
    )
    add_footer(slide)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, prs, "8. 经营与管理成效指标")
    add_bullets(
        slide,
        Inches(0.9),
        Inches(1.9),
        Inches(12.1),
        Inches(4.9),
        [
            "增长：浏览量、咨询量、发布量、活跃用户",
            "转化：咨询转单率、下单率、订金支付率、尾款支付率",
            "履约：合同确认率、变更完成率、结算完成率、平均履约周期",
            "财务：成交额、佣金收入、退款率、放款成功率",
            "服务：会话响应时长、工单关闭时长、SLA 达成率",
            "风控：高风险操作数、告警闭环时长、审计覆盖率",
        ],
        size=20,
    )
    add_footer(slide)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, prs, "9. 当前交付与后续规划")
    add_two_col_cards(
        slide,
        [
            ("当前已具备", "可用于对外汇报、流程演示、运营客服财务协同作业、审计治理落地。"),
            ("当前边界", "主目标为平台内闭环；H5 支付引导回小程序；外部能力分阶段接入。"),
            ("下一步方向", "强化数据运营看板、完善生态接入能力、推进更多自动化作业。"),
            ("汇报结论", "平台已具备“可运营、可治理、可扩展”的甲方交付基础。"),
        ],
    )
    add_footer(slide)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(str(OUTPUT))


if __name__ == "__main__":
    build()
